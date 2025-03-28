from pptx import Presentation

# Create a new presentation
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Machine Learning Web Application"
subtitle.text = "A Streamlined Approach to ML Workflows\nPresented by [Your Name]"

# Slide 2: Introduction
slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction"
content.text = (
    "This project aims to streamline machine learning workflows by creating a user-friendly "
    "web application that supports dataset preprocessing, model selection, and evaluation. "
    "It focuses on simplifying ML tasks for both technical and non-technical users."
)

# Slide 3: Problem Statement
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Problem Statement"
content.text = (
    "Traditional machine learning workflows are often complex and require extensive coding. "
    "This creates a barrier for non-experts and limits the accessibility of ML solutions. "
    "The project addresses these issues by providing an automated and interactive web application."
)

# Slide 4: Literature Survey
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Literature Survey"
content.text = (
    "1. AutoML Frameworks: Automates algorithm and hyperparameter selection using Bayesian optimization.\n"
    "2. Scalable Preprocessing: Uses cloud-based solutions for efficient preprocessing of large datasets.\n"
    "3. Clustering-Based Visualization: Combines PCA and K-Means to improve interpretability of high-dimensional data."
)

# Slide 5: Tools and Technologies
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Tools and Technologies"
content.text = (
    "Python: Programming language for ML algorithms.\n"
    "Streamlit: Framework for building interactive web applications.\n"
    "Scikit-learn: Library for implementing machine learning models.\n"
    "Pandas, NumPy: Tools for data preprocessing and manipulation."
)

# Slide 6: Methodology
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Methodology"
content.text = (
    "1. Dataset Upload: Users upload datasets in CSV format.\n"
    "2. Preprocessing: Missing values, scaling, and encoding handled automatically.\n"
    "3. Model Selection: AutoML selects the best algorithm for the task.\n"
    "4. Training and Evaluation: Models are trained and evaluated on the provided data."
)

# Slide 7: Experiments
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Experiments"
content.text = (
    "Datasets used:\n"
    "- Titanic (Classification)\n"
    "- Wine Quality (Regression)\n"
    "- Mall Customers (Clustering)\n"
    "Results include improved accuracy, reduced mean squared error, and well-defined clusters."
)

# Slide 8: Results
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Results"
content.text = (
    "Classification:\n"
    " - Titanic Dataset: Random Forest achieved 87% accuracy.\n"
    "Regression:\n"
    " - Wine Quality Dataset: Linear Regression achieved R-squared of 0.65.\n"
    "Clustering:\n"
    " - Mall Customers: K-Means achieved a Silhouette Score of 0.75."
)

# Slide 9: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion"
content.text = (
    "The ML web application successfully simplifies complex workflows, making machine learning "
    "more accessible to a wider audience. The integration of AutoML, preprocessing tools, and "
    "interactive visualizations ensures adaptability and scalability for diverse use cases."
)

# Slide 10: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Thank You"
subtitle.text = "Questions? Contact: [Your Email]"

# Save the presentation
prs.save("ML_Web_App_Presentation.pptx")
print("Presentation created successfully!")
